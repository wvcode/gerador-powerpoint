# -*- coding: utf-8 -*-
import json
import os
import shutil
from io import BytesIO

import openai
import requests
import typer
from dotenv import load_dotenv
from langchain import OpenAI
from llama_index import (
    GPTVectorStoreIndex,
    LLMPredictor,
    PromptHelper,
    ServiceContext,
    SimpleDirectoryReader,
)
from pptx import Presentation
from pptx.util import Inches

app = typer.Typer()

load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")


def create_query_engine(source_folder):
    # Carrega lista de documentos do folder que será indexado
    documents = SimpleDirectoryReader(source_folder).load_data()

    # Define parametros para o GPT
    max_input_size = 4096
    num_output = 2048
    max_chunk_overlap = 20

    # Cria o preditor, o prompt helper e o contexto de serviço que permitirá a indexação do conteúdo
    llm_predictor = LLMPredictor(
        llm=OpenAI(temperature=0, model_name="text-davinci-003", max_tokens=num_output)
    )
    prompt_helper = PromptHelper(
        max_input_size=max_input_size,
        num_output=num_output,
        max_chunk_overlap=max_chunk_overlap,
    )
    service_context = ServiceContext.from_defaults(
        llm_predictor=llm_predictor, prompt_helper=prompt_helper
    )

    # Cria o índice baseado nos documentos, utilizando o contexto de serviço
    index = GPTVectorStoreIndex.from_documents(
        documents, service_context=service_context
    )

    # Torna o índice disponível para a 'conversa'
    qe = index.as_query_engine()
    return qe


def download_image_unsplash(keyword):
    url = f"https://api.unsplash.com/photos/random?query={keyword}&client_id={os.getenv('UNSPLASH_ACCESS_KEY')}"
    response = requests.get(url)
    data = response.json()
    image_url = data["urls"]["regular"]
    image_response = requests.get(image_url)
    image_data = BytesIO(image_response.content)
    return image_data


def download_image(description):
    openai.api_key = os.getenv("OPENAI_API_KEY")
    response_data = openai.Image.create(prompt=description, n=2, size="1024x1024")
    image_url = response_data["data"][0]["url"].strip()
    image_response = requests.get(image_url)
    image_data = BytesIO(image_response.content)
    return image_data


def create_slide(prs, title, content, image_concept):
    slide_layout = prs.slide_layouts[1]  # Usando o layout de título e conteúdo
    slide = prs.slides.add_slide(slide_layout)

    # Adicionando título
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Adicionando conteúdo
    content_placeholder = slide.placeholders[1]
    content_placeholder.text = content

    # Adicionando imagem conceito
    image_data = download_image(image_concept)
    image = slide.shapes.add_picture(image_data, Inches(1), Inches(1), height=Inches(3))

    return slide


def create_presentation(source_filename, presentation_filename):
    # Criando uma apresentação vazia
    prs = Presentation()

    slides = json.load(open(source_filename, "r"))

    for slide in slides:
        title1 = slide["title"]
        content1 = slide["content"]
        image_concept1 = slide["image"]
        create_slide(prs, title1, content1, image_concept1)

    # Salvar a apresentação em um arquivo
    prs.save(presentation_filename)


@app.command()
def generate_script_pptx(
    source_folder: str, output_folder: str = "outputs/pptx", index_folder: str = "temp"
):
    files = []
    for root, _, filenames in os.walk(source_folder):
        for filename in filenames:
            if os.path.splitext(filename)[1] == ".pdf":
                files.append(os.path.join(root, filename))

    if len(files) > 0:
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        for item in files:
            print(item)

            if not os.path.exists(index_folder):
                os.makedirs(index_folder)
            else:
                for del_filename in os.listdir(index_folder):
                    os.remove(os.path.join(index_folder, del_filename))

            new_filename = os.path.split(item)[1]
            shutil.copy2(item, os.path.join(index_folder, new_filename))

            if not os.path.exists(
                os.path.join(output_folder, new_filename.replace(".pdf", ".json"))
            ):
                qe = create_query_engine(index_folder)

                comando = """Transforme o conteúdo deste documento em um conjunto de 10 slides
              formatados em um array de objetos json. Cada slide deve ser um elemento do array, com os seguintes campos: 
              title que é o titulo do slide, content que conterá o conteúdo do slide e image, que conterá uma expressão que sumarize o conteúdo do slide, no máximo em 5 palavras.
              O texto deve ser reescrito utilizando um tom técnico."""

                response = qe.query(comando)

                with open(
                    os.path.join(output_folder, new_filename.replace(".pdf", ".json")),
                    "w",
                    encoding="utf8",
                ) as fw:
                    fw.write(str(response))

            if not os.path.exists(
                os.path.join(output_folder, new_filename.replace(".pdf", ".pptx"))
            ):
                create_presentation(
                    os.path.join(output_folder, new_filename.replace(".pdf", ".json")),
                    os.path.join(output_folder, new_filename.replace(".pdf", ".pptx")),
                )


if __name__ == "__main__":
    app()
